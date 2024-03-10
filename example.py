from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Add slides
slide_layout = prs.slide_layouts[1]  # Choose a slide layout (e.g., Title and Content)
slide = prs.slides.add_slide(slide_layout)

# Add title and content to the slide
title = slide.shapes.title
title.text = "Slide 1 Title"

content = slide.placeholders[1]
content.text = "This is slide 1 content."

# Add another slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Slide 2 Title"

content = slide.placeholders[1]
content.text = "This is slide 2 content."

# Save the presentation
prs.save("example.pptx")
