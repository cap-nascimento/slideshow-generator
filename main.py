import pandas as pd
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from datafilter import results



def slide_with_image_and_description(prs, slidetitle, image_path, description):
    slide_layout = prs.slide_layouts[5]  # Choose a layout that fits title and content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = slidetitle
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(135, 60, 49)  # dark red

    left_inch = Inches(2)  # Adjust the left position of the image as needed
    top_inch = Inches(1.5)  # Adjust the top position of the image as needed
    width_inch = Inches(6)  # Adjust the width of the image as needed
    height_inch = Inches(4.5)  # Adjust the height of the image as needed

    slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

    # Add description below the image
    left_inch_desc = left_inch
    top_inch_desc = top_inch + height_inch  # Adjust the spacing between image and description
    width_inch_desc = width_inch
    height_inch_desc = Inches(1)  # Adjust the height of the description text box as needed

    textbox = slide.shapes.add_textbox(left_inch_desc, top_inch_desc, width_inch_desc, height_inch_desc)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = description
    p.font.size = Pt(18)
    text_frame.word_wrap = True

# Example usage:
# slide_with_image_and_description(prs, "Title", "image.jpg", "Description text goes here.")


# Example usage:
# slide_with_image(prs, "Title", "image.jpg")


def simple_slide(prs, slidetitle, contentdata):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = slidetitle
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(135, 60, 49)  # dark red
    # title.text_frame.paragraphs[0].alignment = 1  # Center alignment
    
    content = slide.placeholders[1]
    content.text = contentdata
    content.text_frame.paragraphs[0].font.size = Pt(20)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)  # Dark gray color


def slide_data(prs, slidetitle, data, df_name, max_rows_per_slide=10):
    slide_layout = prs.slide_layouts[5]  # Choose a layout that fits title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = slidetitle
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(135, 60, 49)  # dark red

    rows = len(data[df_name])
    cols = len(data[df_name].columns)

    left_inch = Inches(1)  # Adjust the left position of the table as needed
    top_inch = Inches(1.5)  # Adjust the top position of the table as needed
    width_inch = Inches(4)  # Adjust the width of the table as needed
    height_inch = Inches(4)  # Adjust the height of the table as needed

    total_rows = rows
    num_slides_needed = (total_rows - 1) // max_rows_per_slide + 1

    for slide_num in range(num_slides_needed):
        current_slide = prs.slides[-1]
        table_rows = min(max_rows_per_slide, total_rows - slide_num * max_rows_per_slide)
        slide_table = current_slide.shapes.add_table(table_rows + 1, cols, left_inch, top_inch, width_inch, height_inch).table

        # Set table column widths
        column_widths = column_widths = [Inches(4)] * cols
        for i, width in enumerate(column_widths):
            slide_table.columns[i].width = width

        # Write column headings
        for i, column_name in enumerate(data[df_name].columns):
            slide_table.cell(0, i).text = "Membro" if column_name == "username" else "Acertos"
            slide_table.cell(0, i).text_frame.paragraphs[0].font.bold = True
            slide_table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment

        # Write data to table
        for i, row in enumerate(data[df_name].iloc[slide_num * max_rows_per_slide:slide_num * max_rows_per_slide + table_rows].itertuples(), start=1):
            for j, value in enumerate(row[1:], start=0):
                slide_table.cell(i, j).text = str(value)
                slide_table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment

        if slide_num < num_slides_needed - 1:
            slide = prs.slides.add_slide(slide_layout)

# Example usage:
# slide_data_with_limit(prs, "Title", data, "df_name", max_rows_per_slide=10)


# def slide_data(prs, slidetitle, data, df_name):
#     slide_layout = prs.slide_layouts[5]  # Choose a layout that fits title and content
#     slide = prs.slides.add_slide(slide_layout)

#     title = slide.shapes.title
#     title.text = slidetitle
#     title.text_frame.paragraphs[0].font.bold = True
#     title.text_frame.paragraphs[0].font.size = Pt(32)
#     title.text_frame.paragraphs[0].font.color.rgb = RGBColor(135, 60, 49)  # dark red

#     rows = len(data[df_name]) + 1  # Add 1 for header row
#     cols = len(data[df_name].columns)

#     left_inch = Inches(1)  # Adjust the left position of the table as needed
#     top_inch = Inches(1.5)  # Adjust the top position of the table as needed
#     width_inch = Inches(4)  # Adjust the width of the table as needed
#     height_inch = Inches(4)  # Adjust the height of the table as needed

#     table = slide.shapes.add_table(rows, cols, left_inch, top_inch, width_inch, height_inch).table

#     # Set table column widths
#     # column_widths = [Inches(width_inch / cols)] * cols
#     column_widths = [Inches(4)] * cols
#     for i, width in enumerate(column_widths):
#         table.columns[i].width = width

#     # Write column headings
#     for i, column_name in enumerate(data[df_name].columns):
#         table.cell(0, i).text = "Membro" if column_name == "username" else "Acertos"
#         table.cell(0, i).text_frame.paragraphs[0].font.bold = True
#         table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment

#     # Write data to table
#     for i, row in enumerate(data[df_name].itertuples(), start=1):
#         for j, value in enumerate(row[1:], start=0):
#             table.cell(i, j).text = str(value)
#             table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment



if __name__ == "__main__":
  data = results()
  prs = Presentation()
  
  simple_slide(prs, "Bolin do Oscar", "Resultados")
  simple_slide(prs, "Membros da Academia de Artes e Ciências do Cinema do Paranoá", "")
  
  imgs = {}
  for img in os.listdir('images'):
    imgs[img.split('.')[0]] = 'images/' + img
  for index, row in data["ai_descriptions"].iterrows():
    slide_with_image_and_description(prs, row['username'], imgs[row['username']], row['image'])
  
  slide_data(prs, "Framboesa", data, "razzle")
  slide_data(prs, "Categorias \"Chute\"", data, "lucky")
  slide_data(prs, "Categorias Principais", data, "maincat")
  slide_data(prs, "Categorias Técnicas", data, "technical")
  # slide_data(prs, "Geral", data, "general")
  slide_data(prs, "Ranking Geral", data, "all")
  simple_slide(prs, "Apostas Aleatórias", "")
  for index, row in data["randoms"].iterrows():
    simple_slide(prs, row['username'], row['description'])
  prs.save("results.pptx")
  